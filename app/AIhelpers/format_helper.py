import io
import os
import csv
import fitz
import docx
import pandas as pd
from pptx import Presentation
from PIL import Image
from typing import Iterator, Tuple

from app.ocr import preprocess_for_ocr, safe_ocr

def iterateFilePages(filePath: str) -> Iterator[Tuple[int, str, bool]]:
    ext = os.path.splitext(filePath)[1].lower()

    try:
        if ext == ".pdf":
            yield from extractPdf(filePath)
        elif ext == ".docx":
            yield from extractDocx(filePath)
        elif ext in (".ppt", ".pptx"):
            yield from extractPpt(filePath)
        elif ext in (".xls", ".xlsx"):
            yield from extractExcel(filePath)
        elif ext == ".csv":
            yield from extractCsv(filePath)
        elif ext == ".txt":
            yield from extractTxt(filePath)
        elif ext in (".png", ".jpg", ".jpeg"):
            yield from extractImage(filePath)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as exc:
        # HARD FAILSAFE — ingestion must not crash
        yield 1, f"[INGESTION ERROR]\n{exc}", False

def extractPdf(path: str):
    document = fitz.open(path)
    pageNumber = 1

    for page in document:
        try:
            blocks = page.get_text("blocks")
            texts = [b[4].strip() for b in blocks if b[4].strip()]
            combined = "\n".join(texts)

            # Only trust native text if it's meaningful
            if len(combined) >= 80:
                yield pageNumber, combined, False
            else:
                pix = page.get_pixmap(dpi=200)
                image = Image.frombytes(
                    "RGB",
                    (pix.width, pix.height),
                    pix.samples,
                )
                image = preprocess_for_ocr(image)
                ocrText, used = safe_ocr(image)

                if ocrText:
                    yield pageNumber, "[OCR PAGE]\n" + ocrText, used

        except Exception as exc:
            yield pageNumber, f"[PDF PAGE ERROR]\n{exc}", False

        pageNumber += 1

def extractDocx(path: str):
    document = docx.Document(path)
    buffer = []
    pageNumber = 1

    def flush():
        nonlocal buffer, pageNumber
        if buffer:
            text = "\n".join(buffer)
            buffer = []
            yield pageNumber, text, False
            pageNumber += 1

    for para in document.paragraphs:
        if para.text.strip():
            buffer.append(para.text.strip())
        if len(buffer) >= 6:
            yield from flush()

    yield from flush()

    for rel in document.part._rels.values():
        if "image" in rel.reltype:
            try:
                image = Image.open(io.BytesIO(rel.target_part.blob))
                image = preprocess_for_ocr(image)
                ocrText, used = safe_ocr(image)
                if ocrText:
                    yield pageNumber, "[IMAGE OCR]\n" + ocrText, used
                    pageNumber += 1
            except Exception:
                continue

def extractPpt(path: str):
    prs = Presentation(path)
    slideNumber = 1

    for slide in prs.slides:
        texts = []

        # Slide text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())

        if texts:
            inferred = inferProcessSteps(texts)
            yield slideNumber, "[SLIDE TEXT]\n" + inferred, False

        # Tables
        for shape in slide.shapes:
            if shape.has_table:
                headers = [c.text.strip() for c in shape.table.rows[0].cells]
                rows = []

                for row in shape.table.rows[1:]:
                    row_map = {
                        headers[i]: cell.text.strip()
                        for i, cell in enumerate(row.cells)
                        if i < len(headers) and cell.text.strip()
                    }
                    if row_map:
                        rows.append(row_map)

                if rows:
                    yield slideNumber, "[SLIDE TABLE]\n" + "\n".join(map(str, rows)), False

        # Charts
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                info = [f"Chart type: {chart.chart_type}"]

                if chart.chart_title:
                    info.append(f"Title: {chart.chart_title.text_frame.text}")

                for series in chart.series:
                    values = list(series.values)[:10]
                    info.append(f"{series.name}: {values}")

                yield slideNumber, "[CHART DATA]\n" + "\n".join(info), False

        # Images OCR (diagrams / screenshots)
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                try:
                    image = Image.open(io.BytesIO(shape.image.blob))
                    image = preprocess_for_ocr(image)
                    ocrText, used = safe_ocr(image)
                    if ocrText:
                        yield slideNumber, "[SLIDE IMAGE OCR]\n" + ocrText, used
                except Exception:
                    continue

        slideNumber += 1

def extractExcel(path: str):
    sheets = pd.read_excel(path, sheet_name=None)
    pageNumber = 1

    for name, df in sheets.items():
        try:
            df = df.dropna(how="all")
            if df.empty:
                continue

            rows = df.to_dict(orient="records")
            text = f"[SHEET: {name}]\n" + "\n".join(map(str, rows))

            yield pageNumber, text, False
            pageNumber += 1

        except Exception as exc:
            yield pageNumber, f"[EXCEL ERROR]\n{exc}", False
            pageNumber += 1

def extractCsv(path: str):
    rows = []
    pageNumber = 1

    with open(path, encoding="utf-8", errors="ignore") as file:
        reader = csv.reader(file)
        for row in reader:
            if row:
                rows.append(" | ".join(row))
            if len(rows) >= 15:
                yield pageNumber, "\n".join(rows), False
                rows = []
                pageNumber += 1

    if rows:
        yield pageNumber, "\n".join(rows), False

def extractTxt(path: str):
    with open(path, encoding="utf-8", errors="ignore") as file:
        lines = [l.strip() for l in file if l.strip()]

    buffer = []
    pageNumber = 1

    for line in lines:
        buffer.append(line)
        if len(buffer) >= 10:
            yield pageNumber, "\n".join(buffer), False
            buffer = []
            pageNumber += 1

    if buffer:
        yield pageNumber, "\n".join(buffer), False

def extractImage(path: str):
    try:
        image = Image.open(path)
        image = preprocess_for_ocr(image)
        ocrText, used = safe_ocr(image)

        if ocrText and ocrText.strip():
            yield 1, "[IMAGE OCR]\n" + ocrText.strip(), True
        else:
            yield 1, (
                "This document is an image but OCR could not extract readable text."
            ), True

    except Exception as exc:
        yield 1, f"[IMAGE OCR ERROR]\n{exc}", True

def inferProcessSteps(lines: list[str]) -> str:
    steps = []
    for line in lines:
        lower = line.lower()
        if any(k in lower for k in ("step", "process", "phase", "flow", "stage")):
            steps.append(line)

    return "\n".join(steps) if steps else "\n".join(lines)
